"""
ETL - File Parser Module

Contains format-specific file parsing logic (PDF, DOCX, Excel, PPTX, images, text).
"""

import io
import logging

logger = logging.getLogger(__name__)


async def parse_file_content(content: bytes, filename: str, call_ai_raw) -> tuple[str, dict | None]:
    """
    Parse file content based on filename extension.

    Args:
        content: Raw file bytes
        filename: Original filename (used for extension detection)
        call_ai_raw: Async callable for AI-based OCR

    Returns:
        (extracted_text, error_dict_or_none)
        If error_dict is not None, parsing failed and the dict contains
        the error response to return.
    """
    import asyncio

    text = ""
    fn_lower = filename.lower()

    if fn_lower.endswith(".pdf"):
        from pypdf import PdfReader

        def _parse_pdf():
            pdf_text = ""
            reader = PdfReader(io.BytesIO(content))
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    pdf_text += extracted + "\n"
            return pdf_text

        text = await asyncio.to_thread(_parse_pdf)

    elif fn_lower.endswith((".txt", ".md", ".csv", ".json")):
        text = content.decode("utf-8")

    elif fn_lower.endswith((".png", ".jpg", ".jpeg")):
        import base64

        base64_image = base64.b64encode(content).decode("utf-8")
        payload = {
            "model": "gpt-4o-mini",
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "Transcribe the text in this image accurately. Preserve layout if possible.",
                        },
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"},
                        },
                    ],
                }
            ],
            "max_tokens": 4000,
        }
        try:
            text = await call_ai_raw(payload, endpoint="/chat/completions")
        except Exception as e:
            return "", {"filename": filename, "status": "skipped", "reason": f"OCR Failed: {str(e)}"}

    elif fn_lower.endswith(".docx"):

        def _parse_docx():
            import docx

            doc_obj = docx.Document(io.BytesIO(content))
            return "\n".join([para.text for para in doc_obj.paragraphs])

        try:
            text = await asyncio.to_thread(_parse_docx)
        except Exception as e:
            error_str = str(e)
            if "Bad magic number" in error_str or "File is not a zip file" in error_str:
                reason = "文件格式错误。请确认这是标准的 .docx 文件（OpenXML）。"
            else:
                reason = f"DOCX 解析失败: {error_str}"
            return "", {"filename": filename, "status": "error", "reason": reason}

    elif fn_lower.endswith((".xlsx", ".xls")):

        def _parse_excel():
            try:
                import openpyxl
            except ImportError:
                raise ImportError("openpyxl 未安装，请运行: pip install openpyxl")

            if filename.lower().endswith(".xls"):
                raise ValueError("不支持旧版 .xls 格式（Excel 97-2003）。请将文件另存为 .xlsx 格式后重新上传。")

            wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            parts = []

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                sheet_lines = [f"## 工作表: {sheet_name}\n"]

                rows = list(ws.iter_rows(values_only=True))
                if not rows:
                    sheet_lines.append("（空工作表）\n")
                    parts.append("\n".join(sheet_lines))
                    continue

                for row_idx, row in enumerate(rows):
                    cell_values = [str(cell) if cell is not None else "" for cell in row]
                    row_text = "| " + " | ".join(cell_values) + " |"
                    sheet_lines.append(row_text)
                    if row_idx == 0:
                        separator = "| " + " | ".join(["---"] * len(cell_values)) + " |"
                        sheet_lines.append(separator)

                parts.append("\n".join(sheet_lines))

            wb.close()
            return "\n\n".join(parts)

        try:
            text = await asyncio.to_thread(_parse_excel)
        except (ImportError, ValueError) as e:
            return "", {"filename": filename, "status": "error", "reason": str(e)}
        except Exception as e:
            error_str = str(e)
            if "File is not a zip file" in error_str:
                reason = "文件格式错误。请确认这是标准的 .xlsx 文件（OpenXML）。"
            else:
                reason = f"Excel 解析失败: {error_str}"
            return "", {"filename": filename, "status": "error", "reason": reason}

    elif fn_lower.endswith((".pptx", ".ppt")):

        def _parse_pptx():
            try:
                from pptx import Presentation
            except ImportError:
                raise ImportError("python-pptx 未安装，请运行: pip install python-pptx")

            if filename.lower().endswith(".ppt"):
                raise ValueError(
                    "不支持旧版 .ppt 格式（PowerPoint 97-2003）。请将文件另存为 .pptx 格式后重新上传。"
                )

            prs = Presentation(io.BytesIO(content))
            parts = []

            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_lines = [f"## 幻灯片 {slide_num}"]

                if slide.shapes.title and slide.shapes.title.has_text_frame:
                    slide_lines.append(f"**标题:** {slide.shapes.title.text_frame.text}")

                body_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        if shape == slide.shapes.title:
                            continue
                        for paragraph in shape.text_frame.paragraphs:
                            para_text = paragraph.text.strip()
                            if para_text:
                                body_texts.append(para_text)

                if body_texts:
                    slide_lines.append("\n".join(body_texts))

                parts.append("\n".join(slide_lines))

            return "\n\n".join(parts)

        try:
            text = await asyncio.to_thread(_parse_pptx)
        except (ImportError, ValueError) as e:
            return "", {"filename": filename, "status": "error", "reason": str(e)}
        except Exception as e:
            error_str = str(e)
            if "File is not a zip file" in error_str or "Package not found" in error_str:
                reason = "文件格式错误。请确认这是标准的 .pptx 文件（OpenXML）。"
            else:
                reason = f"PPT 解析失败: {error_str}"
            return "", {"filename": filename, "status": "error", "reason": reason}
    else:
        return "", {"filename": filename, "status": "skipped", "reason": "Unsupported format"}

    return text, None
