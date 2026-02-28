"""
ETL Service - Document processing, extraction, and embedding pipeline
"""

import hashlib
import io
import json
import logging
from typing import Any

import httpx
from pypdf import PdfReader

from app.core.config import settings
from app.core.database import supabase

logger = logging.getLogger(__name__)


class ETLService:
    """
    Enhanced ETL Service using raw HTTP calls (httpx) to maintain
    maximum compatibility with 3rd-party proxies like apiyi.com.
    """

    # Default embedding model (can be overridden by gateway resolution)
    _DEFAULT_EMBEDDING_MODEL = "text-embedding-3-small"

    def __init__(self):
        self.api_key = settings.OPENAI_API_KEY
        # Normalize Base URL: Ensure it ends with /v1
        base_url = settings.AI_BASE_URL if settings.AI_BASE_URL else "https://api.openai.com/v1"
        self.base_url = base_url.rstrip("/")

        # RAG Configurable Parameters (Optimization 5)
        # Defaults: Size=600, Overlap=100
        # Now read from config.py instead of env directly
        from app.core.config import settings as app_settings

        self.chunk_size = app_settings.RAG_CHUNK_SIZE
        self.chunk_overlap = app_settings.RAG_CHUNK_OVERLAP

    async def _call_ai_raw(self, payload: dict, endpoint: str = "/chat/completions") -> str:
        """
        Low-level HTTP call to the AI proxy. Bypass SDK limitations.
        """
        if not self.api_key:
            raise Exception("AI API Key is missing in environment variables")

        # ... logic ...
        headers = {
            "Authorization": f"Bearer {self.api_key}",
            "Content-Type": "application/json",
        }
        url = f"{self.base_url}{endpoint}"

        async with httpx.AsyncClient(timeout=60.0) as client:
            response = await client.post(url, headers=headers, json=payload)
            if response.status_code != 200:
                error_msg = response.text
                logger.error(f"AI Provider Error ({response.status_code}): {error_msg}")
                raise Exception(f"AI provider returned error {response.status_code}")

            data = response.json()
            return data["choices"][0]["message"]["content"]

    def _scrub_pii(self, content: str) -> str:
        """
        Unified PII scrubbing logic.
        Enhances privacy protection for sensitive data before DB storage.
        """
        import re

        # 1. Phone Numbers (Simple 11 digits)
        content = re.sub(r"(?<!\d)1[3-9]\d{9}(?!\d)", "[PHONE_REDACTED]", content)

        # 2. Chinese ID Card (18 digits or 17+X)
        # Matches 18-digit ID cards: 6 (Area) + 8 (DOB) + 4 (Suffix)
        # Mask the DOB part (8 digits) with asterisks
        content = re.sub(r"(?<!\d)(\d{6})\d{8}(\d{3}[\dXx])(?!\d)", r"\1********\2", content)

        # 3. Email Addresses
        content = re.sub(
            r"[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+",
            "[EMAIL_REDACTED]",
            content,
        )

        # 4. API Keys / Secrets (Heuristic)
        content = re.sub(
            r"(?i)(password|passwd|secret|api_key|access_key|token)\s*[:=]\s*[^\s\n,]+",
            r"\1=[SENSITIVE_REDACTED]",
            content,
        )

        # 5. Private Keys
        content = re.sub(
            r"-----BEGIN [A-Z ]+ PRIVATE KEY-----[\s\S]*?-----END [A-Z ]+ PRIVATE KEY-----",
            "[PRIVATE_KEY_REDACTED]",
            content,
        )

        return content

    @staticmethod
    def compute_content_hash(content: bytes) -> str:
        """
        P1 Fix #20: Compute SHA-256 fingerprint of file content for deduplication.
        Identical files will produce the same hash regardless of filename.
        """
        return hashlib.sha256(content).hexdigest()

    async def check_duplicate(
        self, content_hash: str, user_id: str, *, org_id: str | None = None, filename: str | None = None
    ) -> dict | None:
        """
        Multi-level deduplication check:
        1. Exact content hash match (same user) — original P1 Fix #20
        2. Exact content hash match (org-wide) — cross-user dedup
        3. Title similarity match (org-wide) — fuzzy filename dedup
        Returns the existing document record if found, None otherwise.
        """
        if not supabase:
            return None

        try:
            # Level 1: Exact hash match for same user
            res = (
                await supabase.table("documents")
                .select("id, name, status, created_at, owner_id")
                .eq("content_hash", content_hash)
                .eq("owner_id", user_id)
                .limit(1)
                .execute()
            )
            if res.data and len(res.data) > 0:
                return {**res.data[0], "dedup_reason": "exact_hash_same_user"}

            # Level 2: Exact hash match across org
            if org_id:
                org_res = (
                    await supabase.table("documents")
                    .select("id, name, status, created_at, owner_id")
                    .eq("content_hash", content_hash)
                    .eq("organization_id", org_id)
                    .limit(1)
                    .execute()
                )
                if org_res.data and len(org_res.data) > 0:
                    return {**org_res.data[0], "dedup_reason": "exact_hash_org"}

            # Level 3: Title similarity match (org-wide)
            if filename and org_id:
                similar = await self._check_title_similarity(filename, org_id)
                if similar:
                    return similar

        except Exception as e:
            logger.debug(f"Dedup check failed (non-fatal): {e}")
        return None

    async def _check_title_similarity(
        self, filename: str, org_id: str, threshold: float = 0.85
    ) -> dict | None:
        """
        Check for documents with very similar titles within the same org.
        Uses SequenceMatcher for fuzzy matching to catch renamed duplicates.
        """
        from difflib import SequenceMatcher
        import re

        # Normalize filename: strip extension and common prefixes/suffixes
        base_name = filename.rsplit(".", 1)[0] if "." in filename else filename
        normalized = re.sub(r"[\s_\-()（）\[\]【】]+", "", base_name).lower()

        if len(normalized) < 3:
            return None

        try:
            # Fetch recent org documents (limit to avoid full table scan)
            res = (
                await supabase.table("documents")
                .select("id, name, status, created_at, owner_id")
                .eq("organization_id", org_id)
                .in_("status", ["ready", "processing", "pending"])
                .order("created_at", desc=True)
                .limit(200)
                .execute()
            )
            if not res.data:
                return None

            for doc in res.data:
                doc_name = doc.get("name", "")
                doc_base = doc_name.rsplit(".", 1)[0] if "." in doc_name else doc_name
                doc_normalized = re.sub(r"[\s_\-()（）\[\]【】]+", "", doc_base).lower()

                if not doc_normalized:
                    continue

                # Exact normalized match
                if normalized == doc_normalized:
                    return {**doc, "dedup_reason": "title_exact"}

                # Fuzzy match via SequenceMatcher
                ratio = SequenceMatcher(None, normalized, doc_normalized).ratio()
                if ratio >= threshold:
                    return {**doc, "dedup_reason": f"title_similar({ratio:.0%})"}

        except Exception as e:
            logger.debug(f"Title similarity check failed (non-fatal): {e}")
        return None

    async def create_initial_record(
        self,
        filename: str,
        user_id: str,
        status: str = "pending",
        visibility: str = "organization",  # P0 Security Fix #4
        department: str = None,
        content_hash: str = None,  # P1 Fix #20
        category: str = "other",  # H6: Document Category
        organization_id: str = None,  # RAG Fix: org isolation
    ) -> str:
        """
        Creates a placeholder record in the database.

        P0 Security Fix #4: Added visibility parameter
        - 'organization': Visible to all company members (default, for shared knowledge)
        - 'department': Visible only to same department members
        - 'private': Visible only to the uploader (personal notes, drafts)
        """
        if not supabase:
            raise Exception("Supabase not initialized")

        # Validate visibility
        if visibility not in ("private", "department", "organization"):
            visibility = "organization"

        # If department visibility but no department provided, try to get user's department
        if visibility == "department" and not department:
            try:
                user_res = await supabase.table("users").select("department").eq("id", user_id).maybe_single().execute()
                if user_res.data:
                    department = user_res.data.get("department")
            except Exception as e:
                logger.debug(f"Failed to fetch user department: {e}")

        record = {
            "name": filename,
            "status": "pending",
            "progress": 0,
            "stage": "uploading",
            "owner_id": user_id,
            "visibility": visibility,  # P0 Security Fix #4
            "department": department,  # P0 Security Fix #4
            "content_hash": content_hash,  # P1 Fix #20: fingerprint for dedup
            "organization_id": organization_id,  # RAG Fix: org isolation
        }
        res = await supabase.table("documents").insert(record).execute()
        if not res.data:
            raise Exception("Failed to create initial document record")
        return res.data[0]["id"]

    async def _update_progress(self, doc_id: str, progress: int, stage: str, status: str = "processing"):
        """Updates the progress of the document processing."""
        if not doc_id:
            return
        try:
            await supabase.table("documents").update({"progress": progress, "stage": stage, "status": status}).eq(
                "id", doc_id
            ).execute()
        except Exception as e:
            logger.error(f"Failed to update progress for {doc_id}: {e}")

    async def process_file(
        self,
        content: bytes,
        filename: str,
        doc_id: str = None,
        api_key: str = None,
        base_url: str = None,
        user_id: str = None,
        organization_id: str = None,  # RAG Fix: org isolation
    ) -> dict:
        text = ""

        # Initial Progress Update
        await self._update_progress(doc_id, 10, "parsing")

        # Use provided config or fall back to system settings
        # URL Normalization: Extract base even if user provided full endpoint
        raw_url = (base_url or self.base_url).split("/chat/completions")[0].split("/embeddings")[0].rstrip("/")
        if "/v1" not in raw_url and "api.openai.com" not in raw_url:
            active_url = f"{raw_url}/v1" if not raw_url.endswith("/v1") else raw_url
        else:
            active_url = raw_url
        active_key = api_key or self.api_key

        try:
            import asyncio

            # 1. Physical Extraction
            if filename.lower().endswith(".pdf"):
                # Offload CPU-bound task to thread to avoid blocking the event loop (TC-06)
                def _parse_pdf():
                    pdf_text = ""
                    reader = PdfReader(io.BytesIO(content))
                    for page in reader.pages:
                        extracted = page.extract_text()
                        if extracted:
                            pdf_text += extracted + "\n"
                    return pdf_text

                text = await asyncio.to_thread(_parse_pdf)
            elif filename.lower().endswith((".txt", ".md", ".csv", ".json")):
                text = content.decode("utf-8")
            elif filename.lower().endswith((".png", ".jpg", ".jpeg")):
                # OCR is already calling an external AI API (async), so it's fine.
                import base64

                base64_image = base64.b64encode(content).decode("utf-8")
                payload = {
                    "model": "gpt-4o-mini",
                    "messages": [
                        {
                            "role": "user",
                            "content": [
                                {
                                    "type": "text",
                                    "text": "Transcribe the text in this image accurately. Preserve layout if possible.",
                                },
                                {
                                    "type": "image_url",
                                    "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"},
                                },
                            ],
                        }
                    ],
                    "max_tokens": 4000,
                }
                try:
                    text = await self._call_ai_raw(payload, endpoint="/chat/completions")
                except Exception as e:
                    return {
                        "filename": filename,
                        "status": "skipped",
                        "reason": f"OCR Failed: {str(e)}",
                    }
            elif filename.lower().endswith(".docx"):

                def _parse_docx():
                    import docx

                    doc_obj = docx.Document(io.BytesIO(content))
                    return "\n".join([para.text for para in doc_obj.paragraphs])

                try:
                    text = await asyncio.to_thread(_parse_docx)
                except Exception as e:
                    error_str = str(e)
                    if "Bad magic number" in error_str or "File is not a zip file" in error_str:
                        return {
                            "filename": filename,
                            "status": "error",
                            "reason": "文件格式错误。请确认这是标准的 .docx 文件（OpenXML）。",
                        }
                    return {
                        "filename": filename,
                        "status": "error",
                        "reason": f"DOCX 解析失败: {error_str}",
                    }
            elif filename.lower().endswith((".xlsx", ".xls")):

                def _parse_excel():
                    try:
                        import openpyxl
                    except ImportError:
                        raise ImportError("openpyxl 未安装，请运行: pip install openpyxl")

                    if filename.lower().endswith(".xls"):
                        raise ValueError("不支持旧版 .xls 格式（Excel 97-2003）。请将文件另存为 .xlsx 格式后重新上传。")

                    wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
                    parts = []

                    for sheet_name in wb.sheetnames:
                        ws = wb[sheet_name]
                        sheet_lines = [f"## 工作表: {sheet_name}\n"]

                        rows = list(ws.iter_rows(values_only=True))
                        if not rows:
                            sheet_lines.append("（空工作表）\n")
                            parts.append("\n".join(sheet_lines))
                            continue

                        # Build markdown table format for better structure preservation
                        for row_idx, row in enumerate(rows):
                            cell_values = [str(cell) if cell is not None else "" for cell in row]
                            row_text = "| " + " | ".join(cell_values) + " |"
                            sheet_lines.append(row_text)
                            # Add header separator after first row
                            if row_idx == 0:
                                separator = "| " + " | ".join(["---"] * len(cell_values)) + " |"
                                sheet_lines.append(separator)

                        parts.append("\n".join(sheet_lines))

                    wb.close()
                    return "\n\n".join(parts)

                try:
                    text = await asyncio.to_thread(_parse_excel)
                except ImportError as e:
                    return {
                        "filename": filename,
                        "status": "error",
                        "reason": str(e),
                    }
                except ValueError as e:
                    return {
                        "filename": filename,
                        "status": "error",
                        "reason": str(e),
                    }
                except Exception as e:
                    error_str = str(e)
                    if "File is not a zip file" in error_str:
                        return {
                            "filename": filename,
                            "status": "error",
                            "reason": "文件格式错误。请确认这是标准的 .xlsx 文件（OpenXML）。",
                        }
                    return {
                        "filename": filename,
                        "status": "error",
                        "reason": f"Excel 解析失败: {error_str}",
                    }
            elif filename.lower().endswith((".pptx", ".ppt")):

                def _parse_pptx():
                    try:
                        from pptx import Presentation
                    except ImportError:
                        raise ImportError("python-pptx 未安装，请运行: pip install python-pptx")

                    if filename.lower().endswith(".ppt"):
                        raise ValueError(
                            "不支持旧版 .ppt 格式（PowerPoint 97-2003）。请将文件另存为 .pptx 格式后重新上传。"
                        )

                    prs = Presentation(io.BytesIO(content))
                    parts = []

                    for slide_num, slide in enumerate(prs.slides, start=1):
                        slide_lines = [f"## 幻灯片 {slide_num}"]

                        # Extract slide title
                        if slide.shapes.title and slide.shapes.title.has_text_frame:
                            slide_lines.append(f"**标题:** {slide.shapes.title.text_frame.text}")

                        # Extract text from all text frames
                        body_texts = []
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                # Skip the title shape (already extracted above)
                                if shape == slide.shapes.title:
                                    continue
                                for paragraph in shape.text_frame.paragraphs:
                                    para_text = paragraph.text.strip()
                                    if para_text:
                                        body_texts.append(para_text)

                        if body_texts:
                            slide_lines.append("\n".join(body_texts))

                        # Only add slide if it has any text content
                        if len(slide_lines) > 1:
                            parts.append("\n".join(slide_lines))

                    return "\n\n".join(parts)

                try:
                    text = await asyncio.to_thread(_parse_pptx)
                except ImportError as e:
                    return {
                        "filename": filename,
                        "status": "error",
                        "reason": str(e),
                    }
                except ValueError as e:
                    return {
                        "filename": filename,
                        "status": "error",
                        "reason": str(e),
                    }
                except Exception as e:
                    error_str = str(e)
                    if "File is not a zip file" in error_str or "Package not found" in error_str:
                        return {
                            "filename": filename,
                            "status": "error",
                            "reason": "文件格式错误。请确认这是标准的 .pptx 文件（OpenXML）。",
                        }
                    return {
                        "filename": filename,
                        "status": "error",
                        "reason": f"PPT 解析失败: {error_str}",
                    }
            else:
                return {
                    "filename": filename,
                    "status": "skipped",
                    "reason": "Unsupported format",
                }

            if not text.strip():
                return {
                    "filename": filename,
                    "status": "error",
                    "reason": "No text content found",
                }

            # Update Progress: Extraction Done
            await self._update_progress(doc_id, 30, "analyzing")

            # 2. Sequential Processing
            success, details = await self.extract_metadata_via_ai(text, filename, active_key, active_url)

            # Update Progress: Metadata Done
            await self._update_progress(doc_id, 70, "embedding")

            if success:
                try:
                    # Save Logic: If doc_id exists, update it. If not, create new (legacy path)
                    if doc_id:
                        # Scrub PII
                        safe_text = self._scrub_pii(text)

                        details["full_text_context"] = safe_text[:100000]

                        await supabase.table("documents").update(
                            {
                                "extracted_data": details,
                                "doc_type": details.get("doc_type", "other"),
                                "status": "processing",  # Still processing embeddings
                            }
                        ).eq("id", doc_id).execute()
                    else:
                        # Legacy creation if no doc_id passed
                        doc_id = await self._save_to_db(
                            filename,
                            details,
                            text,
                            user_id=user_id,
                            status="processing",
                            organization_id=organization_id,
                        )

                    # P0 Security Fix: Scrub PII before generating embeddings
                    # This prevents sensitive data (phone numbers, ID cards, emails)
                    # from being vectorized and retrievable via semantic search
                    safe_text_for_embedding = self._scrub_pii(text)

                    # Generate embeddings with PII-scrubbed text
                    embedding_success = await self._generate_embeddings(
                        safe_text_for_embedding,
                        doc_id,
                        filename,
                        active_key,
                        active_url,
                        organization_id=organization_id,
                    )

                    if embedding_success:
                        # Finalize status
                        await self._update_progress(doc_id, 100, "completed", status="ready")
                        return {
                            "filename": filename,
                            "status": "success",
                            "document_id": doc_id,
                            "metadata": details,
                        }
                    else:
                        # P1: Rollback/Mark failed if embeddings fail
                        await supabase.table("documents").update(
                            {
                                "status": "failed",
                                "error_log": "Embedding generation partially failed",
                            }
                        ).eq("id", doc_id).execute()
                        return {
                            "filename": filename,
                            "status": "partial_success",
                            "reason": "文档已记录，但向量索引失败，搜索可能受限。",
                        }

                except Exception as db_err:
                    logger.error(f"DB Error: {db_err}")
                    if doc_id:
                        await supabase.table("documents").update({"status": "error", "error_log": str(db_err)}).eq(
                            "id", doc_id
                        ).execute()
                    return {
                        "filename": filename,
                        "status": "error",
                        "reason": f"数据库写入失败: {str(db_err)}",
                    }
            else:
                return {
                    "filename": filename,
                    "status": "error",
                    "reason": f"AI 解析失败: {details.get('error')}",
                }

        except Exception as e:
            logger.error(f"ETL Panic: {str(e)}")
            if doc_id:
                await self._update_progress(doc_id, 0, "failed", status="error")
            return {
                "filename": filename,
                "status": "error",
                "reason": f"系统崩溃: {str(e)}",
            }

    async def extract_metadata_via_ai(
        self, text: str, filename: str, api_key: str, base_url: str
    ) -> tuple[bool, dict[str, Any]]:
        """
        Uses AI to extract structured metadata (JSON) from raw text.
        Supports Tender Analysis (Redlines, Deviations) for 'bid' type documents.
        """
        # Filename-based pre-classification hint for AI
        _fn_lower = filename.lower()
        _filename_hint = "other"
        _filename_keywords = {
            "product": [
                "彩页",
                "产品资料",
                "产品手册",
                "产品说明",
                "规格书",
                "datasheet",
                "brochure",
                "产品目录",
                "catalog",
            ],
            "contract": ["合同", "协议", "contract"],
            "tender": ["招标", "招标文件", "tender", "招标公告", "采购需求"],
            "bid": ["投标", "投标书", "标书", "bid", "应标", "响应文件"],
            "proposal": ["方案", "proposal"],
            "invoice": ["发票", "invoice"],
        }
        for _dtype, _keywords in _filename_keywords.items():
            if any(kw in _fn_lower for kw in _keywords):
                _filename_hint = _dtype
                break

        prompt = f"""
        # Role (角色设定)
        你是一位拥有 20 年经验的资深招投标专家 (Senior Bid Manager) 和项目经理。你擅长快速阅读复杂的招标文件，精准捕捉关键信息，分析潜在风险，并制定高胜率的投标策略。

        # Task (任务)
        请分析以下文件内容（全文上下文），按 **6 个核心模块** 进行结构化提取和分析。

        文件名: {filename}
        文件名预分类提示: {_filename_hint}（仅供参考，请根据实际内容判断）
        文件内容片段:
        {text}

        # Output Format (输出格式)
        请输出两个独立的部分，严禁将 Markdown 报告包含在 JSON 字段中。

        第一部分：元数据 (JSON)
        [METADATA_JSON]
        {{
            "doc_type": "tender" | "bid" | "contract" | "product" | "proposal" | "invoice" | "other",
            "client_name": "采购方/客户名称",
            "amount": 预算金额(数字)或null,
            "date": "YYYY-MM-DD",
            "tags": ["核心标签"],
            "redlines": ["提取模块2中的核心否决项(简练列表)"],
            "technical_deviations": ["提取模块3/6中的技术风险点(简练列表)"]
        }}
        [/METADATA_JSON]

        doc_type 分类说明：
        - "tender": 招标文件、采购需求文件（客户/甲方发布的需求文档，包含技术要求、资质条件、评分标准等）
        - "bid": 投标文件、应标书、响应文件（我方编写的投标响应文档，包含技术方案、报价、资质证明等）
        - "contract": 合同、协议
        - "product": 产品资料、产品彩页、产品手册、规格书、产品说明书、产品目录
        - "proposal": 解决方案、项目方案、技术方案
        - "invoice": 发票、账单
        - "other": 其他无法归类的文件

        第二部分：完整分析报告 (Markdown)
        [ANALYSIS_REPORT]
        ## 模块 1：项目概况与时间表
        ...
        ## 模块 6：风险预警与专家建议
        ...
        [/ANALYSIS_REPORT]
        """

        payload = {
            "messages": [
                {
                    "role": "system",
                    "content": "You are a senior Tender Analyst. Output structured data.",
                },
                {"role": "user", "content": prompt},
            ],
            "temperature": 0.2,
        }

        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
        }

        async def call_ai_model(model_name: str, retries=1) -> tuple[bool, Any]:
            """Helper to call AI with retry logic"""
            payload["model"] = model_name
            try:
                logger.info(f"Attempting AI Analysis with model: {model_name}...")
                async with httpx.AsyncClient(timeout=90.0) as client:
                    response = await client.post(f"{base_url}/chat/completions", headers=headers, json=payload)

                    if response.status_code != 200:
                        logger.warning(f"Model {model_name} failed: {response.status_code} - {response.text}")
                        return False, None

                    return True, response.json()
            except Exception as e:
                logger.warning(f"Model {model_name} processing error: {str(e)}")
                return False, None

        # 1. Try Primary Model (Bleeding Edge)
        success, response_json = await call_ai_model("gemini-3-pro-preview")

        # 2. Fallback to Stable Model if primary fails
        if not success:
            logger.warning("Primary model failed. Falling back to Gemini-2.5-Pro...")
            success, response_json = await call_ai_model("gemini-2.5-pro")

        if not success or not response_json:
            return False, {"error": "All AI models failed to process the document."}

        try:
            content = response_json["choices"][0]["message"]["content"]

            # Extract JSON
            import re

            json_match = re.search(r"\[METADATA_JSON\](.*?)\[/METADATA_JSON\]", content, re.DOTALL)
            report_match = re.search(r"\[ANALYSIS_REPORT\](.*?)\[/ANALYSIS_REPORT\]", content, re.DOTALL)

            metadata = {}
            if json_match:
                try:
                    metadata = json.loads(json_match.group(1).strip())
                except json.JSONDecodeError as e:
                    logger.warning(f"JSON Parse Failed: {e}")

            if report_match:
                metadata["full_analysis_markdown"] = report_match.group(1).strip()
            elif not json_match:
                # Fallback: if no blocks found, try raw JSON parse (if model ignored instructions)
                try:
                    clean_json = content.replace("```json", "").replace("```", "").strip()
                    metadata = json.loads(clean_json)
                except json.JSONDecodeError:
                    pass  # Fallback parsing also failed

            if not metadata:
                raise Exception("Failed to parse AI output format")

            return True, metadata
        except Exception as e:

            logger.warning(f"Metadata Extraction Failed: {e}")
            # Fallback metadata
            return True, {
                "doc_type": "other",
                "summary": "AI 解析失败，请手动审阅。",
                "client_name": None,
                "amount": 0,
                "date": None,
                "tags": ["解析失败"],
                "redlines": [],
                "technical_deviations": [],
            }

    async def _save_to_db(
        self,
        filename: str,
        metadata: dict,
        text: str = "",
        user_id: str = None,
        status: str = "ready",
        visibility: str = "organization",  # P0 Security Fix #4
        department: str = None,
        organization_id: str = None,  # RAG Fix: org isolation
    ) -> str:
        """
        Save document to database with visibility control.

        P0 Security Fix #4: Added visibility parameter for three-tier access control.
        """
        if not supabase:
            raise Exception("Supabase not initialized")

        # Validate visibility
        if visibility not in ("private", "department", "organization"):
            visibility = "organization"

        safe_text = self._scrub_pii(text)
        metadata["full_text_context"] = safe_text[:100000]

        record = {
            "name": filename,
            "doc_type": metadata.get("doc_type", "other"),
            "extracted_data": metadata,
            "version": 1,
            "owner_id": user_id,
            "status": status,
            "visibility": visibility,  # P0 Security Fix #4
            "department": department,  # P0 Security Fix #4
            "organization_id": organization_id,  # RAG Fix: org isolation
        }
        res = await supabase.table("documents").insert(record).execute()
        if not res.data:
            raise Exception("Empty response from database")
        return res.data[0]["id"]

    async def _generate_embeddings(
        self,
        text: str,
        doc_id: str,
        filename: str,
        api_key: str,
        base_url: str,
        organization_id: str = None,
    ) -> bool:
        """
        Batch Embeddings with partial success tracking.
        """
        batch_size = 50
        current_batch_text = []
        all_success = True

        # Resolve embedding model dynamically via gateway
        embedding_model = self._DEFAULT_EMBEDDING_MODEL
        active_api_key = api_key
        active_base_url = base_url
        try:
            from app.services.llm_helpers import resolve_embedding_config

            emb_config = await resolve_embedding_config(organization_id or "default")
            if emb_config.get("model"):
                embedding_model = emb_config["model"]
            if emb_config.get("api_key"):
                active_api_key = emb_config["api_key"]
            if emb_config.get("base_url"):
                active_base_url = emb_config["base_url"].rstrip("/")
                if "/v1" not in active_base_url and "api.openai.com" not in active_base_url:
                    active_base_url = f"{active_base_url}/v1"
        except Exception:
            pass

        async def _process_batch(batch_texts):
            try:
                payload = {"model": embedding_model, "input": batch_texts}
                headers = {"Authorization": f"Bearer {active_api_key}"}
                async with httpx.AsyncClient(timeout=45.0) as client:
                    resp = await client.post(f"{active_base_url}/embeddings", headers=headers, json=payload)
                    if resp.status_code == 200:
                        embeddings_data = resp.json()["data"]
                        records = []
                        for i, item in enumerate(embeddings_data):
                            records.append(
                                {
                                    "document_id": doc_id,
                                    "content": batch_texts[i],
                                    "embedding": item["embedding"],
                                    "metadata": {"source": filename},
                                    "organization_id": organization_id,
                                }
                            )
                        await supabase.table("document_embeddings").insert(records).execute()
                        return True
                    return False
            except Exception as e:
                logger.error(f"Batch embedding failed: {e}")
                return False

        # Use new dynamic size and overlap
        for chunk in self._semantic_chunk(text, size=self.chunk_size, overlap=self.chunk_overlap):
            current_batch_text.append(chunk)
            if len(current_batch_text) >= batch_size:
                if not await _process_batch(current_batch_text):
                    all_success = False
                current_batch_text = []

        if current_batch_text and not await _process_batch(current_batch_text):
            all_success = False

        return all_success

    def _semantic_chunk(self, text: str, size: int = 600, overlap: int = 100):
        """
        P2 Fix: Improved chunking strategy.
        Tries to split by double newlines (paragraphs) first, then falls back
        to sliding window if paragraphs are too large.
        """
        # 1. Clean up excessive whitespace
        import re

        text = re.sub(r"\n{3,}", "\n\n", text)

        # 2. Initial split by double newlines
        paragraphs = text.split("\n\n")

        current_chunk = ""

        for p in paragraphs:
            # If paragraph itself is too large, split it by sentences or characters
            if len(p) > size:
                # If we have something in current_chunk, yield it
                if current_chunk:
                    yield current_chunk
                    current_chunk = ""

                # Split large paragraph by sliding window
                start = 0
                while start < len(p):
                    end = start + size
                    chunk = p[start:end]
                    yield chunk
                    start += size - overlap
            else:
                # If current_chunk + new paragraph is within limit
                if len(current_chunk) + len(p) < size:
                    current_chunk += ("\n\n" if current_chunk else "") + p
                else:
                    # Yield current and start new
                    if current_chunk:
                        yield current_chunk
                    current_chunk = p

        if current_chunk:
            yield current_chunk


etl_service = ETLService()
